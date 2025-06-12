import os
import fitz  # PyMuPDF for PDFs
import openai
import tempfile
from flask import Flask, request, jsonify
from pptx import Presentation
from werkzeug.utils import secure_filename

# 🔐 Set your OpenAI API Key here or via environment variable
openai.api_key = os.getenv("OPENAI_API_KEY")

app = Flask(__name__)
ALLOWED_EXTENSIONS = {'pdf', 'pptx'}

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pdf(file_path):
    doc = fitz.open(file_path)
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_ppt(file_path):
    prs = Presentation(file_path)
    return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])

def summarize_with_gpt(text):
    prompt = f"Summarize the following content:\n\n{text[:4000]}"  # truncate to stay within GPT limit
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        max_tokens=300
    )
    return response['choices'][0]['message']['content']

@app.route('/summarize', methods=['POST'])
def summarize_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file part in the request"}), 400

    file = request.files['file']
    if file.filename == '' or not allowed_file(file.filename):
        return jsonify({"error": "Invalid or missing file"}), 400

    filename = secure_filename(file.filename)
    with tempfile.NamedTemporaryFile(delete=False, suffix=filename) as temp_file:
        file.save(temp_file.name)
        ext = filename.rsplit('.', 1)[1].lower()

        try:
            if ext == 'pdf':
                text = extract_text_from_pdf(temp_file.name)
            elif ext == 'pptx':
                text = extract_text_from_ppt(temp_file.name)
            else:
                return jsonify({"error": "Unsupported file type"}), 400
        finally:
            os.remove(temp_file.name)

    summary = summarize_with_gpt(text)
    return jsonify({"summary": summary})

if __name__ == '__main__':
    app.run(debug=True)

