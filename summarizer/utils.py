import fitz  # PyMuPDF
from pptx import Presentation
import openai
import os

openai.api_key = os.getenv("OPENAI_API_KEY")

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_ppt(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def summarize_text(text):
    prompt = f"Summarize the following content:\n{text}"
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": "You are a helpful assistant that summarizes documents."},
            {"role": "user", "content": prompt},
        ],
        temperature=0.5,
        max_tokens=500
    )
    return response["choices"][0]["message"]["content"]
